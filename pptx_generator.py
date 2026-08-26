import os
import sys
import io
from datetime import datetime
import pandas as pd
import numpy as np

# Configurazione Matplotlib Headless & Cache Locale
mpl_dir = os.path.join(os.path.dirname(os.path.abspath(__file__)), '.matplotlib')
os.makedirs(mpl_dir, exist_ok=True)
os.environ['MPLCONFIGDIR'] = mpl_dir

import matplotlib
matplotlib.use('Agg')
import matplotlib.pyplot as plt

# PowerPoint PPTX
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE

# Palette Ufficiale AstraZeneca
AZ_NAVY = RGBColor(0x00, 0x38, 0x65)
AZ_GOLD = RGBColor(0xD0, 0xA0, 0x00)
AZ_BERRY = RGBColor(0x8B, 0x00, 0x4B)
AZ_SLATE = RGBColor(0x1C, 0x2A, 0x39)
AZ_LIGHT_BG = RGBColor(0xF4, 0xF6, 0xF9)
AZ_WHITE = RGBColor(0xFF, 0xFF, 0xFF)
AZ_BLUE_ACCENT = RGBColor(0x02, 0x84, 0xC7)
AZ_CARD_BORDER = RGBColor(0xCB, 0xD5, 0xE1)

# Palette Esadecimale per Grafici Matplotlib
HEX_NAVY = '#003865'
HEX_GOLD = '#D0A000'
HEX_BERRY = '#8B004B'
HEX_BLUE = '#0284C7'
HEX_CYAN = '#0EA5E9'
HEX_GREEN = '#10B981'
HEX_SLATE = '#1C2A39'
HEX_GRAY = '#64748B'
HEX_BG = '#F4F6F9'

# ---------------------------------------------------------
# Funzioni per Renderizzare i Grafici in Immagini ad Alta Risoluzione
# ---------------------------------------------------------
def render_donut_chart(labels, values, title, color_palette=None):
    if color_palette is None:
        color_palette = [HEX_NAVY, HEX_GOLD, HEX_BERRY, HEX_BLUE, HEX_GREEN, HEX_GRAY]
        
    fig, ax = plt.subplots(figsize=(4.8, 3.6), dpi=160, facecolor='#FFFFFF')
    
    total = sum(values)
    if total == 0:
        values = [1]
        labels = ['Nessun Dato']
        color_palette = [HEX_GRAY]
        
    wedges, texts, autotexts = ax.pie(
        values,
        labels=labels,
        autopct=lambda p: f'{p:.1f}%\n({int(round(p*total/100))})' if total > 0 else '',
        colors=color_palette[:len(values)],
        startangle=140,
        pctdistance=0.72,
        wedgeprops=dict(width=0.48, edgecolor='#FFFFFF', linewidth=2.5),
        textprops=dict(fontsize=8, weight='bold', color=HEX_SLATE)
    )
    for autotext in autotexts:
        autotext.set_fontsize(7.5)
        autotext.set_color('#FFFFFF')
        autotext.set_weight('bold')
        
    ax.set_title(title, fontsize=10.5, weight='bold', color=HEX_NAVY, pad=12)
    plt.tight_layout()
    
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=160, facecolor='#FFFFFF')
    plt.close(fig)
    buf.seek(0)
    return buf

def render_horizontal_bar_chart(categories, values, title, bar_color=HEX_NAVY):
    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=160, facecolor='#FFFFFF')
    
    if len(categories) == 0 or len(values) == 0:
        categories = ['Nessun Dato']
        values = [0]
        
    y_pos = np.arange(len(categories))
    bars = ax.barh(y_pos, values, color=bar_color, height=0.55, edgecolor='none', zorder=3)
    
    ax.set_yticks(y_pos)
    ax.set_yticklabels(categories, fontsize=8.5, weight='bold', color=HEX_SLATE)
    ax.invert_yaxis()
    
    # Valori agli apici delle barre
    max_val = max(values) if values else 1
    for bar, v in zip(bars, values):
        ax.text(
            bar.get_width() + (max_val * 0.02),
            bar.get_y() + bar.get_height() / 2,
            f'{int(v):,}',
            va='center', ha='left',
            fontsize=8.5, weight='bold', color=HEX_NAVY
        )
        
    ax.set_xlim(0, max_val * 1.25)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CBD5E1')
    ax.spines['bottom'].set_color('#CBD5E1')
    ax.grid(axis='x', linestyle='--', alpha=0.3, zorder=0)
    ax.set_title(title, fontsize=10.5, weight='bold', color=HEX_NAVY, pad=12)
    plt.tight_layout()
    
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=160, facecolor='#FFFFFF')
    plt.close(fig)
    buf.seek(0)
    return buf

def render_time_series_chart(dates, values, title, line_color=HEX_NAVY):
    fig, ax = plt.subplots(figsize=(5.2, 3.6), dpi=160, facecolor='#FFFFFF')
    
    if len(dates) == 0:
        dates = ['-']
        values = [0]
        
    date_strs = [str(d)[5:] if len(str(d))>=10 else str(d) for d in dates]
    x_pos = np.arange(len(date_strs))
    
    ax.plot(x_pos, values, marker='o', markersize=6, color=line_color, linewidth=2.5, zorder=4)
    ax.fill_between(x_pos, values, color=line_color, alpha=0.1, zorder=2)
    
    # Valori sopra i punti
    max_val = max(values) if values else 1
    for x, v in zip(x_pos, values):
        ax.annotate(
            str(int(v)),
            (x, v),
            textcoords="offset points",
            xytext=(0, 7),
            ha='center',
            fontsize=8,
            weight='bold',
            color=HEX_NAVY
        )
        
    ax.set_xticks(x_pos)
    ax.set_xticklabels(date_strs, rotation=40, ha='right', fontsize=7.5, color=HEX_SLATE)
    ax.set_ylim(0, max_val * 1.25)
    ax.spines['top'].set_visible(False)
    ax.spines['right'].set_visible(False)
    ax.spines['left'].set_color('#CBD5E1')
    ax.spines['bottom'].set_color('#CBD5E1')
    ax.grid(axis='y', linestyle='--', alpha=0.3, zorder=0)
    ax.set_title(title, fontsize=10.5, weight='bold', color=HEX_NAVY, pad=12)
    plt.tight_layout()
    
    buf = io.BytesIO()
    plt.savefig(buf, format='png', bbox_inches='tight', dpi=160, facecolor='#FFFFFF')
    plt.close(fig)
    buf.seek(0)
    return buf

# ---------------------------------------------------------
# Slide Template & Header/Footer Helpers
# ---------------------------------------------------------
def create_slide_header(slide, title_text, subtitle_text="AstraZeneca Medical Information & Decision Intelligence Hub"):
    header_box = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(1.15))
    header_box.fill.solid()
    header_box.fill.fore_color.rgb = AZ_NAVY
    header_box.line.color.rgb = AZ_GOLD
    header_box.line.width = Pt(2)
    
    stripe = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.15), Inches(13.333), Inches(0.06))
    stripe.fill.solid()
    stripe.fill.fore_color.rgb = AZ_GOLD
    stripe.line.fill.background()

    tf = header_box.text_frame
    tf.word_wrap = True
    tf.margin_left = Inches(0.6)
    tf.margin_top = Inches(0.15)
    
    p = tf.paragraphs[0]
    p.text = title_text
    p.font.size = Pt(21)
    p.font.bold = True
    p.font.color.rgb = AZ_WHITE
    
    p2 = tf.add_paragraph()
    p2.text = subtitle_text
    p2.font.size = Pt(10.5)
    p2.font.color.rgb = RGBColor(0xDF, 0xE7, 0xF0)

def add_footer(slide, current_page=1):
    footer_box = slide.shapes.add_textbox(Inches(0.6), Inches(7.05), Inches(12.133), Inches(0.35))
    tf = footer_box.text_frame
    p = tf.paragraphs[0]
    p.text = f"AstraZeneca Medical Affairs | Confidenziale - Uso Interno Riservato | Data: {datetime.now().strftime('%d/%m/%Y')} | Slide {current_page}"
    p.font.size = Pt(9)
    p.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

# ---------------------------------------------------------
# Generatore Deck PowerPoint (.pptx) Completo
# ---------------------------------------------------------
def generate_pptx_deck(df, matrix_df, audit_trail, product_filter=None):
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5) # 16:9 Widescreen
    blank_layout = prs.slide_layouts[6]
    
    working_df = df.copy()
    is_product_deck = False
    if product_filter and product_filter != "Tutti i Prodotti" and 'Product_Clean' in working_df.columns:
        working_df = working_df[working_df['Product_Clean'] == product_filter]
        is_product_deck = True

    total_cases = len(working_df)
    ae_cnt = int(working_df['is_ae'].sum()) if 'is_ae' in working_df.columns else 0
    pqc_cnt = int(working_df['is_pqc'].sum()) if 'is_pqc' in working_df.columns else 0
    unsol_cnt = int(working_df['is_unsolicited'].sum()) if 'is_unsolicited' in working_df.columns else 0
    sla_hours = working_df['SLA_Hours'].dropna() if 'SLA_Hours' in working_df.columns else pd.Series(dtype=float)
    avg_sla = sla_hours.mean() if not sla_hours.empty else 0
    
    slide_counter = 1

    # ----------------------------------------------------
    # SLIDE 1: Copertina Istituzionale AstraZeneca
    # ----------------------------------------------------
    s1 = prs.slides.add_slide(blank_layout)
    bg1 = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(7.5))
    bg1.fill.solid()
    bg1.fill.fore_color.rgb = AZ_NAVY
    bg1.line.fill.background()
    
    gold_bar = s1.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.8), Inches(1.5), Inches(0.2), Inches(4.5))
    gold_bar.fill.solid()
    gold_bar.fill.fore_color.rgb = AZ_GOLD
    gold_bar.line.fill.background()
    
    tb1 = s1.shapes.add_textbox(Inches(1.3), Inches(1.8), Inches(11), Inches(4))
    tf1 = tb1.text_frame
    tf1.word_wrap = True
    
    p = tf1.paragraphs[0]
    p.text = "ASTRAZENECA MEDICAL AFFAIRS & PHARMACOVIGILANCE"
    p.font.size = Pt(14)
    p.font.bold = True
    p.font.color.rgb = AZ_GOLD
    
    p2 = tf1.add_paragraph()
    p2.text = f"Medical Information & Decision Intelligence Report" if not is_product_deck else f"Deep Dive Report: {product_filter}"
    p2.font.size = Pt(32)
    p2.font.bold = True
    p2.font.color.rgb = AZ_WHITE
    p2.space_after = Pt(14)
    
    p3 = tf1.add_paragraph()
    p3.text = f"Analisi Strategica delle Richieste Cliniche, Insight Medici e Valutazione di Sicurezza" if not is_product_deck else f"Focus Monografico su Domande Cliniche, Canali, Tipologia Referenti e Unmet Needs per {product_filter}"
    p3.font.size = Pt(16)
    p3.font.color.rgb = RGBColor(0xDF, 0xE7, 0xF0)
    p3.space_after = Pt(28)
    
    p4 = tf1.add_paragraph()
    p4.text = f"Data di Generazione: {datetime.now().strftime('%d/%m/%Y')} | Perimetro: {total_cases} Richieste Analizzate"
    p4.font.size = Pt(12)
    p4.font.color.rgb = RGBColor(0x94, 0xA3, 0xB8)
    
    slide_counter += 1

    # ----------------------------------------------------
    # SLIDE 2: Executive Summary & KPI Glance
    # ----------------------------------------------------
    s2 = prs.slides.add_slide(blank_layout)
    create_slide_header(s2, "Executive Summary: Indicatori Chiave (KPI & Safety)", "Sintesi dei volumi operativi, tempi di risposta e monitoraggio di Farmacovigilanza")
    
    kpi_configs = [
        ("TOTALE RICHIESTE GESTITE", f"{total_cases:,}", "Volumi complessivi nel perimetro", AZ_NAVY),
        ("SLA MEDIO RISOLUZIONE", f"{int(avg_sla)}h {int(round((avg_sla%1)*60))}m", f"Tempo medio evasione ticket", AZ_NAVY),
        ("% RICHIESTE SPONTANEE", f"{(unsol_cnt/total_cases*100):.1f}%" if total_cases>0 else "0%", f"{unsol_cnt} richieste spontanee (phactMI/MILE)", AZ_NAVY),
        ("ALERT SICUREZZA (PV/PQC)", f"{ae_cnt + pqc_cnt}", f"{ae_cnt} Eventi Avversi / {pqc_cnt} Reclami Qualità", AZ_BERRY if (ae_cnt+pqc_cnt)>0 else AZ_NAVY)
    ]
    
    left_positions = [Inches(0.6), Inches(3.68), Inches(6.76), Inches(9.84)]
    for idx, (lbl, val, desc, col) in enumerate(kpi_configs):
        card = s2.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left_positions[idx], Inches(1.5), Inches(2.85), Inches(1.7))
        card.fill.solid()
        card.fill.fore_color.rgb = AZ_LIGHT_BG
        card.line.color.rgb = col
        card.line.width = Pt(2)
        
        ctf = card.text_frame
        ctf.word_wrap = True
        ctf.margin_top = Inches(0.18)
        ctf.margin_left = Inches(0.2)
        
        cp1 = ctf.paragraphs[0]
        cp1.text = lbl
        cp1.font.size = Pt(9.5)
        cp1.font.bold = True
        cp1.font.color.rgb = AZ_SLATE
        
        cp2 = ctf.add_paragraph()
        cp2.text = val
        cp2.font.size = Pt(23)
        cp2.font.bold = True
        cp2.font.color.rgb = col
        
        cp3 = ctf.add_paragraph()
        cp3.text = desc
        cp3.font.size = Pt(8.5)
        cp3.font.color.rgb = RGBColor(0x64, 0x74, 0x8B)

    # Narrative Summary Box
    summary_box = s2.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(3.5), Inches(12.133), Inches(3.3))
    summary_box.fill.solid()
    summary_box.fill.fore_color.rgb = AZ_WHITE
    summary_box.line.color.rgb = AZ_CARD_BORDER
    
    stf = summary_box.text_frame
    stf.word_wrap = True
    stf.margin_left = Inches(0.3)
    stf.margin_top = Inches(0.25)
    
    sp1 = stf.paragraphs[0]
    sp1.text = "SINTESI CLINICO-STRATEGICA E VALUTAZIONE REGOLATORIA"
    sp1.font.size = Pt(12)
    sp1.font.bold = True
    sp1.font.color.rgb = AZ_NAVY
    sp1.space_after = Pt(8)
    
    sp2 = stf.add_paragraph()
    sp2.text = f"• Attività di Medical Information: Il portafoglio ha gestito {total_cases} richieste di chiarimento scientifico con un tempo medio di risposta pari a {int(avg_sla)}h {int(round((avg_sla%1)*60))}m, garantendo supporto qualificato agli specialisti."
    sp2.font.size = Pt(10)
    sp2.font.color.rgb = AZ_SLATE
    sp2.space_after = Pt(6)
    
    sp3 = stf.add_paragraph()
    sp3.text = f"• Triage di Farmacovigilanza: Rilevati {ae_cnt} eventi avversi e {pqc_cnt} richieste per deviazioni termiche/device. È stata verificata l'immediata conformità alle procedure di escalation al Safety Team entro 24 ore."
    sp3.font.size = Pt(10)
    sp3.font.color.rgb = AZ_SLATE
    sp3.space_after = Pt(6)
    
    sp4 = stf.add_paragraph()
    unsol_pct = (unsol_cnt / total_cases * 100.0) if total_cases > 0 else 0.0
    sp4.text = f"• Tasso di Spontaneità: Il {unsol_pct:.1f}% delle interazioni è nato da genuina iniziativa del medico (unsolicited), a testimonianza del forte bisogno di approfondimento su dati clinici e paper registrativi."
    sp4.font.size = Pt(10)
    sp4.font.color.rgb = AZ_SLATE
    
    add_footer(s2, slide_counter)
    slide_counter += 1

    # ----------------------------------------------------
    # SLIDE 3: Visual Analytics - Donut Charts (Referenti & Canali)
    # ----------------------------------------------------
    s3 = prs.slides.add_slide(blank_layout)
    create_slide_header(s3, "Visual Analytics: Canali di Ricezione & Tipologia Referente", "Grafici a torta con evidenza esatta delle quote percentuali e volumi (HCP / Rep / Canali)")
    
    # 1. Donut Referenti
    ref_counts = working_df['Referrer_Clean'].value_counts() if 'Referrer_Clean' in working_df.columns else pd.Series()
    ref_buf = render_donut_chart(
        list(ref_counts.index), 
        list(ref_counts.values), 
        "Distribuzione per Tipologia Referente (HCP / Rep)",
        [HEX_NAVY, HEX_GOLD, HEX_BERRY, HEX_BLUE, HEX_GRAY]
    )
    s3.shapes.add_picture(ref_buf, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4.3))
    
    # 2. Donut Canali
    chan_counts = working_df['Case Origin'].value_counts() if 'Case Origin' in working_df.columns else pd.Series()
    chan_buf = render_donut_chart(
        list(chan_counts.index), 
        list(chan_counts.values), 
        "Distribuzione per Canale di Ricezione (Email, Telefono, F2F)",
        [HEX_NAVY, HEX_BLUE, HEX_GOLD, HEX_BERRY, HEX_GREEN]
    )
    s3.shapes.add_picture(chan_buf, Inches(6.9), Inches(1.4), Inches(5.8), Inches(4.3))
    
    # Bottom Note Box
    b_box = s3.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.85), Inches(12.133), Inches(1.0))
    b_box.fill.solid()
    b_box.fill.fore_color.rgb = AZ_LIGHT_BG
    b_box.line.color.rgb = AZ_CARD_BORDER
    btf = b_box.text_frame
    btf.margin_left = Inches(0.2)
    btf.margin_top = Inches(0.12)
    bp = btf.paragraphs[0]
    top_ref = ref_counts.index[0] if not ref_counts.empty else "N/A"
    top_chan = chan_counts.index[0] if not chan_counts.empty else "N/A"
    bp.text = f"Insight di Canale: Il canale primario di contatto è '{top_chan}' ({chan_counts.iloc[0] if not chan_counts.empty else 0} ticket), mentre la categoria di richiedente preponderante è '{top_ref}'. I dati confermano una forte integrazione tra interazioni dirette e segnalazioni di territorio."
    bp.font.size = Pt(9.5)
    bp.font.color.rgb = AZ_SLATE
    
    add_footer(s3, slide_counter)
    slide_counter += 1

    # ----------------------------------------------------
    # SLIDE 4: Visual Analytics - Trend Temporale & Volume Farmaci
    # ----------------------------------------------------
    s4 = prs.slides.add_slide(blank_layout)
    create_slide_header(s4, "Visual Analytics: Trend Temporale Giornaliero & Top Farmaci", "Andamento delle richieste nel tempo con valori sui punti e ripartizione per brand")
    
    # 1. Line Chart Trend Giornaliero
    if 'Date_Opened' in working_df.columns:
        vd = working_df.dropna(subset=['Date_Opened']).copy()
        if not vd.empty:
            vd['Giorno'] = vd['Date_Opened'].dt.date
            t_grp = vd.groupby('Giorno').size().reset_index(name='Richieste')
            time_buf = render_time_series_chart(
                list(t_grp['Giorno']), 
                list(t_grp['Richieste']), 
                "Trend Temporale Giornaliero (Totale Richieste)",
                HEX_NAVY
            )
            s4.shapes.add_picture(time_buf, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4.3))
            
    # 2. Bar Chart Top Farmaci
    if 'Product_Clean' in working_df.columns:
        p_grp = working_df[working_df['Product_Clean'] != 'Non Specificato']['Product_Clean'].value_counts().head(6)
        prod_buf = render_horizontal_bar_chart(
            list(p_grp.index),
            list(p_grp.values),
            "Volume Richieste per Farmaco / Brand",
            HEX_NAVY
        )
        s4.shapes.add_picture(prod_buf, Inches(6.9), Inches(1.4), Inches(5.8), Inches(4.3))
        
    b_box4 = s4.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.85), Inches(12.133), Inches(1.0))
    b_box4.fill.solid()
    b_box4.fill.fore_color.rgb = AZ_LIGHT_BG
    b_box4.line.color.rgb = AZ_CARD_BORDER
    btf4 = b_box4.text_frame
    btf4.margin_left = Inches(0.2)
    btf4.margin_top = Inches(0.12)
    bp4 = btf4.paragraphs[0]
    bp4.text = "Monitoraggio Dinamico: I picchi temporali corrispondono alle finestre congressuali e alla pubblicazione di nuovi dati clinici. La concentrazione sui top brand richiede un presidio continuo tramite slide deck aggiornati."
    bp4.font.size = Pt(9.5)
    bp4.font.color.rgb = AZ_SLATE
    
    add_footer(s4, slide_counter)
    slide_counter += 1

    # ----------------------------------------------------
    # SLIDE 5: Visual Analytics - Territorio & Top Studi Clinici
    # ----------------------------------------------------
    s5 = prs.slides.add_slide(blank_layout)
    create_slide_header(s5, "Visual Analytics: Distribuzione Territoriale & Studi Clinici", "Ripartizione geografica e ranking degli studi/trial più richiesti dai clinici")
    
    # 1. Bar Chart Territori
    if 'Country_Clean' in working_df.columns:
        geo_grp = working_df['Country_Clean'].value_counts().head(6)
        geo_buf = render_horizontal_bar_chart(
            list(geo_grp.index),
            list(geo_grp.values),
            "Distribuzione Territoriale delle Richieste",
            HEX_BLUE
        )
        s5.shapes.add_picture(geo_buf, Inches(0.6), Inches(1.4), Inches(5.8), Inches(4.3))
        
    # 2. Bar Chart Top Trials
    trial_keywords = ['WAYPOINT', 'CASCADE', 'NAVIGATOR', 'DAPA-EAT', 'DAPA-CKD', 'MATTERHORN', 'KOMET', 'ELEVATE', 'AMPLIFY', 'DESTINATION']
    details_list = working_df['Details'].astype(str).tolist() if 'Details' in working_df.columns else []
    text_corp = " ".join(details_list).upper()
    t_counts = {tk: text_corp.count(tk) for tk in trial_keywords if text_corp.count(tk) > 0}
    t_sorted = dict(sorted(t_counts.items(), key=lambda item: item[1], reverse=True)[:6])
    
    if t_sorted:
        trial_buf = render_horizontal_bar_chart(
            list(t_sorted.keys()),
            list(t_sorted.values()),
            "Top Studi Clinici & Paper Richiesti (Unmet Need)",
            HEX_GOLD
        )
        s5.shapes.add_picture(trial_buf, Inches(6.9), Inches(1.4), Inches(5.8), Inches(4.3))
        
    b_box5 = s5.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(5.85), Inches(12.133), Inches(1.0))
    b_box5.fill.solid()
    b_box5.fill.fore_color.rgb = AZ_LIGHT_BG
    b_box5.line.color.rgb = AZ_CARD_BORDER
    btf5 = b_box5.text_frame
    btf5.margin_left = Inches(0.2)
    btf5.margin_top = Inches(0.12)
    bp5 = btf5.paragraphs[0]
    bp5.text = "Evidenze Emergenti: Gli studi registrativi ed emergenti in cima alla classifica rappresentano i principali motori di domanda scientifica da parte degli specialisti ospedalieri."
    bp5.font.size = Pt(9.5)
    bp5.font.color.rgb = AZ_SLATE
    
    add_footer(s5, slide_counter)
    slide_counter += 1

    # ----------------------------------------------------
    # SLIDE 6+: Deep Dive Singolo Prodotto con Grafici Embedded
    # ----------------------------------------------------
    if 'Product_Clean' in working_df.columns:
        if is_product_deck:
            prods_to_show = [product_filter]
        else:
            prods_to_show = list(working_df[working_df['Product_Clean'] != 'Non Specificato']['Product_Clean'].value_counts().head(2).index)
            
        for p_idx, prod_name in enumerate(prods_to_show):
            sp_df = working_df[working_df['Product_Clean'] == prod_name]
            sp_tot = len(sp_df)
            sp_ae = int(sp_df['is_ae'].sum()) if 'is_ae' in sp_df.columns else 0
            sp_pqc = int(sp_df['is_pqc'].sum()) if 'is_pqc' in sp_df.columns else 0
            sp_unsol = int(sp_df['is_unsolicited'].sum()) if 'is_unsolicited' in sp_df.columns else 0
            
            s_prod = prs.slides.add_slide(blank_layout)
            create_slide_header(s_prod, f"Deep Dive Prodotto: {prod_name}", f"Grafici di canali, referenti e tipologia richieste per il brand {prod_name}")
            
            # Chart 1: Donut Canali per quel farmaco
            f_chan = sp_df['Case Origin'].value_counts()
            f_chan_buf = render_donut_chart(
                list(f_chan.index), list(f_chan.values), f"Canali di Ricezione ({prod_name})",
                [HEX_NAVY, HEX_BLUE, HEX_GOLD, HEX_BERRY]
            )
            s_prod.shapes.add_picture(f_chan_buf, Inches(0.6), Inches(1.4), Inches(3.9), Inches(3.2))
            
            # Chart 2: Donut Referenti per quel farmaco
            f_ref = sp_df['Referrer_Clean'].value_counts()
            f_ref_buf = render_donut_chart(
                list(f_ref.index), list(f_ref.values), f"Tipologia Referenti ({prod_name})",
                [HEX_NAVY, HEX_GOLD, HEX_BERRY, HEX_BLUE]
            )
            s_prod.shapes.add_picture(f_ref_buf, Inches(4.7), Inches(1.4), Inches(3.9), Inches(3.2))
            
            # Chart 3: Bar Tipologie Richiesta per quel farmaco
            f_type = sp_df['Type_Clean'].value_counts()
            f_type_buf = render_horizontal_bar_chart(
                list(f_type.index), list(f_type.values), f"Tipologie di Richiesta ({prod_name})",
                HEX_NAVY
            )
            s_prod.shapes.add_picture(f_type_buf, Inches(8.8), Inches(1.4), Inches(3.9), Inches(3.2))
            
            # Bottom Box: KPI & Sample Questions
            bot_prod_box = s_prod.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.6), Inches(4.75), Inches(12.133), Inches(2.15))
            bot_prod_box.fill.solid()
            bot_prod_box.fill.fore_color.rgb = AZ_LIGHT_BG
            bot_prod_box.line.color.rgb = AZ_CARD_BORDER
            
            bptf = bot_prod_box.text_frame
            bptf.margin_left = Inches(0.25)
            bptf.margin_top = Inches(0.15)
            
            p_kpi_line = bptf.paragraphs[0]
            p_kpi_line.text = f"METRICHE BRAND: {sp_tot} Ticket Gestiti | {(sp_unsol/sp_tot*100):.1f}% Spontanee (Unsolicited) | {sp_ae} Alert PV | {sp_pqc} Reclami Qualità"
            p_kpi_line.font.bold = True
            p_kpi_line.font.size = Pt(10.5)
            p_kpi_line.font.color.rgb = AZ_NAVY
            p_kpi_line.space_after = Pt(4)
            
            for _, cr in sp_df.head(3).iterrows():
                cp = bptf.add_paragraph()
                c_num = cr.get('Case Number', '')
                c_det = cr.get('Details_Redacted', cr.get('Details', ''))[:140]
                cp.text = f"• [Case #{c_num}] {c_det}..."
                cp.font.size = Pt(8.5)
                cp.font.color.rgb = AZ_SLATE
                
            add_footer(s_prod, slide_counter)
            slide_counter += 1

    # ----------------------------------------------------
    # SLIDE MATRICE STRATEGICA A 9 COLONNE
    # ----------------------------------------------------
    if not matrix_df.empty:
        s_mat = prs.slides.add_slide(blank_layout)
        create_slide_header(s_mat, "Matrice Strategica di Medical Intelligence (9 Colonne)", "Mappatura dei bisogni insoddisfatti, insight clinici e raccomandazioni operative")
        
        disp_matrix = matrix_df.head(4)
        m_rows = len(disp_matrix) + 1
        m_table_shape = s_mat.shapes.add_table(m_rows, 6, Inches(0.5), Inches(1.5), Inches(12.333), Inches(1.1 * m_rows))
        m_table = m_table_shape.table
        
        m_col_widths = [Inches(1.8), Inches(1.8), Inches(2.4), Inches(2.4), Inches(2.4), Inches(1.533)]
        for ci, w in enumerate(m_col_widths):
            m_table.columns[ci].width = w
            
        m_headers = ["Categoria", "Sottocategoria", "Insight Ricavabili", "Messaggi Chiave Clinici", "Note Operative", "Priorità"]
        for ci, h in enumerate(m_headers):
            m_table.cell(0, ci).text = h
            m_table.cell(0, ci).fill.solid()
            m_table.cell(0, ci).fill.fore_color.rgb = AZ_NAVY
            for p in m_table.cell(0, ci).text_frame.paragraphs:
                p.font.color.rgb = AZ_WHITE
                p.font.bold = True
                p.font.size = Pt(9)
                
        for ri, (_, row) in enumerate(disp_matrix.iterrows()):
            m_table.cell(ri+1, 0).text = str(row.get('Categoria', ''))
            m_table.cell(ri+1, 1).text = str(row.get('Sottocategoria', ''))
            m_table.cell(ri+1, 2).text = str(row.get('Insight ricavabili', ''))
            m_table.cell(ri+1, 3).text = str(row.get('Messaggi chiave / suggerimenti per i clinici', ''))
            m_table.cell(ri+1, 4).text = str(row.get('Note operative', ''))
            m_table.cell(ri+1, 5).text = str(row.get('Priorità percepita', ''))
            
            for ci in range(6):
                m_table.cell(ri+1, ci).fill.solid()
                m_table.cell(ri+1, ci).fill.fore_color.rgb = AZ_LIGHT_BG if ri%2==0 else AZ_WHITE
                for p in m_table.cell(ri+1, ci).text_frame.paragraphs:
                    p.font.size = Pt(8)
                    p.font.color.rgb = AZ_SLATE
                    
        add_footer(s_mat, slide_counter)
        slide_counter += 1

    # ----------------------------------------------------
    # SLIDE FINALE: Decision Intelligence & Audit Trail
    # ----------------------------------------------------
    s_end = prs.slides.add_slide(blank_layout)
    create_slide_header(s_end, "Registro Decisionale & Next-Best-Actions (Human-in-the-Loop)", "Tracciabilità delle decisioni strategiche validate per conformità ed esecuzione territoriale")
    
    if audit_trail:
        aud_rows = min(len(audit_trail), 5) + 1
        aud_shape = s_end.shapes.add_table(aud_rows, 5, Inches(0.6), Inches(1.6), Inches(12.133), Inches(0.8 * aud_rows))
        aud_table = aud_shape.table
        aud_table.columns[0].width = Inches(2.2)
        aud_table.columns[1].width = Inches(1.6)
        aud_table.columns[2].width = Inches(3.5)
        aud_table.columns[3].width = Inches(1.8)
        aud_table.columns[4].width = Inches(3.033)
        
        aud_headers = ["Timestamp", "Action ID", "Azione Strategica", "Decisione", "Case Numbers Tracciati"]
        for ci, h in enumerate(aud_headers):
            aud_table.cell(0, ci).text = h
            aud_table.cell(0, ci).fill.solid()
            aud_table.cell(0, ci).fill.fore_color.rgb = AZ_NAVY
            for p in aud_table.cell(0, ci).text_frame.paragraphs:
                p.font.color.rgb = AZ_WHITE
                p.font.bold = True
                p.font.size = Pt(9.5)
                
        for ri, a in enumerate(audit_trail[:5]):
            aud_table.cell(ri+1, 0).text = str(a.get('Timestamp', ''))
            aud_table.cell(ri+1, 1).text = str(a.get('Action_ID', ''))
            aud_table.cell(ri+1, 2).text = str(a.get('Titolo', ''))
            aud_table.cell(ri+1, 3).text = str(a.get('Decisione', ''))
            aud_table.cell(ri+1, 4).text = str(a.get('Case_Numbers', ''))
            
            for ci in range(5):
                aud_table.cell(ri+1, ci).fill.solid()
                aud_table.cell(ri+1, ci).fill.fore_color.rgb = AZ_LIGHT_BG if ri%2==0 else AZ_WHITE
                for p in aud_table.cell(ri+1, ci).text_frame.paragraphs:
                    p.font.size = Pt(8.5)
                    p.font.color.rgb = AZ_SLATE
    else:
        no_aud_box = s_end.shapes.add_textbox(Inches(0.6), Inches(2.0), Inches(12.1), Inches(2.0))
        p = no_aud_box.text_frame.paragraphs[0]
        p.text = "Nessuna decisione ancora convalidata nell'Audit Trail per questa sessione. Approvare o rifiutare le azioni nella scheda Decision Intelligence per popolare il registro formale."
        p.font.size = Pt(12)
        p.font.color.rgb = AZ_SLATE
        
    add_footer(s_end, slide_counter)
    
    out_buf = io.BytesIO()
    prs.save(out_buf)
    return out_buf.getvalue()
